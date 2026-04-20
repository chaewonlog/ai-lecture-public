from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ══════════════════════════════════════════════
#  BOLD TYPOGRAPHY v2 — Refined Design System
# ══════════════════════════════════════════════

# Palette — layered blacks + blue spectrum
BG_BLACK     = RGBColor(0x0D, 0x0D, 0x0D)
CARD_BLACK   = RGBColor(0x16, 0x16, 0x16)
GHOST        = RGBColor(0x1F, 0x1F, 0x1F)   # for giant bg numbers on dark
GHOST_LIGHT  = RGBColor(0xEC, 0xEC, 0xEC)   # for giant bg numbers on white
OFF_BLACK    = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY    = RGBColor(0x3A, 0x3A, 0x3A)
MID_GRAY     = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY   = RGBColor(0xBB, 0xBB, 0xBB)
SUBTLE_BG    = RGBColor(0xF5, 0xF5, 0xF5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

# Blue spectrum
BLUE_DEEP    = RGBColor(0x00, 0x44, 0xCC)
BLUE         = RGBColor(0x00, 0x66, 0xFF)
BLUE_LIGHT   = RGBColor(0x4D, 0x94, 0xFF)
BLUE_PALE    = RGBColor(0xE6, 0xEF, 0xFF)
RED_SOFT     = RGBColor(0xE0, 0x40, 0x40)

FONT = "Pretendard"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
ML = Inches(1.2)
CW = Inches(10.9)

# ── Core helpers ──
def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def txt(slide, left, top, width, height, text,
        size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
        line_spacing=1.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    p.space_after = Pt(0)
    pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return box

def multi(slide, left, top, width, height, lines,
          size=18, color=DARK_GRAY, spacing=12, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.color.rgb = line[1] if len(line) > 1 else color
            p.font.bold = line[2] if len(line) > 2 else bold
        else:
            p.text = line
            p.font.color.rgb = color
            p.font.bold = bold
        p.font.size = Pt(size)
        p.font.name = FONT
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
    return box

def line(slide, left, top, width, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.035))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def vline(slide, left, top, height, color=LIGHT_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.015), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def page(slide, n, total):
    txt(slide, W - Inches(1.5), H - Inches(0.5), Inches(1.2), Inches(0.3),
        f"{n} / {total}", size=9, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def ghost_num(slide, text, color=GHOST, right_align=False):
    """Giant background number for visual impact"""
    x = W - Inches(5.5) if right_align else Inches(-0.5)
    txt(slide, x, Inches(0.5), Inches(6), Inches(4),
        text, size=200, bold=True, color=color, align=PP_ALIGN.RIGHT if right_align else PP_ALIGN.LEFT)

def deco_circle(slide, x, y, size, color):
    """Decorative circle"""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(size), Inches(size))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()

def deco_dot(slide, x, y, color=BLUE):
    """Small accent dot"""
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.12), Inches(0.12))
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()

def card(slide, left, top, width, height, color=SUBTLE_BG):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    return c

def table_clean(slide, rows, cols, data, left, top, width, row_h=Inches(0.55)):
    height = row_h * rows
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for i in range(rows):
        for j in range(cols):
            cell = tbl.cell(i, j)
            cell.text = data[i][j]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_GRAY
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = OFF_BLACK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else SUBTLE_BG
    tbl_xml = tbl._tbl
    for cell_elem in tbl_xml.iter(qn('a:tcPr')):
        for bn in ['lnL', 'lnR', 'lnT', 'lnB']:
            b = cell_elem.find(qn(f'a:{bn}'))
            if b is not None:
                b.set('w', '0')
    return tbl


TOTAL = 22

# ═══════════════════════════════════════════
# 1 — TITLE
# ═══════════════════════════════════════════
s = blank(); bg(s, BG_BLACK)
deco_circle(s, W - Inches(5), Inches(-2), 7, GHOST)
deco_circle(s, W - Inches(3), Inches(4.5), 4, RGBColor(0x14, 0x14, 0x14))
deco_dot(s, Inches(1.2), Inches(1.5), BLUE)
deco_dot(s, Inches(1.65), Inches(1.5), BLUE_LIGHT)
deco_dot(s, Inches(2.1), Inches(1.5), RGBColor(0x33, 0x33, 0x33))

txt(s, ML, Inches(2.2), Inches(11), Inches(1.3),
    "AI 잘 쓰는 사람은", size=54, bold=True, color=WHITE, line_spacing=1.05)
txt(s, ML, Inches(3.5), Inches(11), Inches(1.3),
    "뭐가 다를까.", size=54, bold=True, color=BLUE, line_spacing=1.05)
line(s, ML, Inches(5.1), Inches(2.5))
txt(s, ML, Inches(5.5), Inches(8), Inches(0.5),
    "오늘 배우고 내일 바로 쓰는 AI 활용법", size=18, color=MID_GRAY)
page(s, 1, TOTAL)

# ═══════════════════════════════════════════
# 2 — 손들기
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
ghost_num(s, "?", GHOST_LIGHT, right_align=True)
deco_dot(s, ML, Inches(0.8), BLUE)

txt(s, ML, Inches(1.0), CW, Inches(1.8),
    "잠깐,\n손 한번\n들어볼까요?", size=50, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(3.6), Inches(2))

y = Inches(4.2)
for text, clr, sz in [
    ("AI 써보신 분?", LIGHT_GRAY, 26),
    ("매일 쓰시는 분?", MID_GRAY, 30),
    ("결과에 만족하시는 분?", OFF_BLACK, 36),
]:
    deco_dot(s, ML, y + Inches(0.15), BLUE if clr == OFF_BLACK else RGBColor(0xDD, 0xDD, 0xDD))
    txt(s, Inches(1.7), y, Inches(8), Inches(0.6), text, size=sz, bold=(clr == OFF_BLACK), color=clr)
    y += Inches(0.75)
page(s, 2, TOTAL)

# ═══════════════════════════════════════════
# 3 — 같은 도구, 다른 결과
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
ghost_num(s, "vs", GHOST_LIGHT, right_align=True)
deco_dot(s, ML, Inches(0.8), BLUE)

txt(s, ML, Inches(1.0), CW, Inches(1.2),
    "같은 도구,\n다른 결과.", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(2.8), Inches(2))

# big number comparison
txt(s, ML, Inches(3.5), Inches(5), Inches(0.4),
    "블로그 초안", size=14, bold=True, color=MID_GRAY)
txt(s, ML, Inches(3.9), Inches(2.5), Inches(0.8),
    "3시간", size=36, bold=True, color=LIGHT_GRAY)
txt(s, Inches(3.5), Inches(4.0), Inches(1), Inches(0.5),
    "→", size=28, color=BLUE)
txt(s, Inches(4.5), Inches(3.9), Inches(2.5), Inches(0.8),
    "30분", size=36, bold=True, color=BLUE)

txt(s, ML, Inches(5.0), Inches(5), Inches(0.4),
    "회계감사보고서 분석", size=14, bold=True, color=MID_GRAY)
txt(s, ML, Inches(5.4), Inches(2.5), Inches(0.8),
    "3~5일", size=36, bold=True, color=LIGHT_GRAY)
txt(s, Inches(3.5), Inches(5.5), Inches(1), Inches(0.5),
    "→", size=28, color=BLUE)
txt(s, Inches(4.5), Inches(5.4), Inches(2.5), Inches(0.8),
    "2시간", size=36, bold=True, color=BLUE)

# right side quote
txt(s, Inches(8.0), Inches(4.0), Inches(4.5), Inches(2.0),
    "차이는\n도구가 아니라\n방법입니다.", size=30, bold=True, color=OFF_BLACK, line_spacing=1.3)
line(s, Inches(8.0), Inches(3.8), Inches(1.5))
page(s, 3, TOTAL)

# ═══════════════════════════════════════════
# 4 — 오늘 알려드릴 3가지
# ═══════════════════════════════════════════
s = blank(); bg(s, BG_BLACK)
deco_circle(s, Inches(-2), Inches(-2), 5, GHOST)
deco_circle(s, W - Inches(3), Inches(5), 3, RGBColor(0x14, 0x14, 0x14))

txt(s, ML, Inches(0.8), CW, Inches(1.0),
    "오늘 알려드릴", size=20, color=MID_GRAY)
txt(s, ML, Inches(1.2), CW, Inches(1.0),
    "3가지.", size=60, bold=True, color=WHITE)
line(s, ML, Inches(2.5), Inches(2))

y = Inches(3.3)
for num, title in [("01", "AI한테 나를 기억시켜라"),
                   ("02", "잘 된 거 뜯어서 내 걸로 만들어라"),
                   ("03", "나만의 봇을 만들어라")]:
    deco_dot(s, ML, y + Inches(0.12), BLUE)
    txt(s, Inches(1.7), y, Inches(1.2), Inches(0.5),
        num, size=13, bold=True, color=BLUE_LIGHT)
    txt(s, Inches(2.8), y, Inches(8), Inches(0.5),
        title, size=26, color=WHITE)
    if num != "03":
        line(s, Inches(2.8), y + Inches(0.6), Inches(7), GHOST)
    y += Inches(0.9)
page(s, 4, TOTAL)

# ═══════════════════════════════════════════
# 5 — Section ①
# ═══════════════════════════════════════════
s = blank(); bg(s, BG_BLACK)
ghost_num(s, "01", GHOST, right_align=True)
deco_circle(s, Inches(8), Inches(3), 6, RGBColor(0x12, 0x12, 0x12))
deco_dot(s, ML, Inches(1.2), BLUE)

txt(s, ML, Inches(1.4), Inches(2), Inches(0.4),
    "PART 01", size=13, bold=True, color=BLUE_LIGHT)
txt(s, ML, Inches(2.2), CW, Inches(2.5),
    "AI한테 나를\n기억시켜라.", size=60, bold=True, color=WHITE, line_spacing=1.08)
line(s, ML, Inches(4.8), Inches(2.5))
txt(s, ML, Inches(5.2), Inches(10), Inches(0.5),
    "매번 설명하는 사람 vs 한 번만 설명하는 사람", size=20, color=MID_GRAY)
page(s, 5, TOTAL)

# ═══════════════════════════════════════════
# 6 — 어떤 AI?
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.8), BLUE)
txt(s, ML, Inches(1.0), CW, Inches(1.2),
    "어떤 AI를\n써야 할까?", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
txt(s, ML, Inches(2.8), Inches(10), Inches(0.5),
    "AI는 핸드폰 기종 고르는 것과 같아요. 써봐야 내 손에 맞는 게 뭔지 알아요.",
    size=15, color=MID_GRAY)
line(s, ML, Inches(3.4), Inches(2))

# 3 cards
col_w = Inches(3.3)
gap = Inches(0.35)
items = [
    ("ChatGPT", "범용, 가장 대중적", BLUE),
    ("Claude", "글쓰기, 분석에 강함", BLUE_LIGHT),
    ("Gemini", "구글 생태계 연동\nNotebookLM으로 자료 정리", BLUE_DEEP),
]
for i, (name, desc, dot_c) in enumerate(items):
    x = ML + (col_w + gap) * i
    card(s, x, Inches(3.9), col_w, Inches(2.8), SUBTLE_BG)
    deco_dot(s, x + Inches(0.4), Inches(4.3), dot_c)
    txt(s, x + Inches(0.7), Inches(4.15), col_w - Inches(1), Inches(0.5),
        name, size=26, bold=True, color=OFF_BLACK)
    line(s, x + Inches(0.4), Inches(4.9), Inches(1), dot_c)
    txt(s, x + Inches(0.4), Inches(5.2), col_w - Inches(0.8), Inches(1.2),
        desc, size=14, color=MID_GRAY, line_spacing=1.5)

txt(s, ML, Inches(7.0), Inches(10), Inches(0.4),
    "자신에게 잘 맞는 걸 테스트해보고 고르세요.", size=16, bold=True, color=BLUE)
page(s, 6, TOTAL)

# ═══════════════════════════════════════════
# 7 — 무료로 쓰는 법
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
ghost_num(s, "FREE", GHOST_LIGHT, right_align=True)
deco_dot(s, ML, Inches(0.8), BLUE)

txt(s, ML, Inches(1.0), CW, Inches(1.2),
    "유료 AI,\n무료로 쓰는 법.", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(2.8), Inches(2))

# Two cards
card(s, ML, Inches(3.5), Inches(5.2), Inches(3.3), SUBTLE_BG)
deco_dot(s, Inches(1.6), Inches(3.9), BLUE)
txt(s, Inches(2.0), Inches(3.8), Inches(4), Inches(0.5),
    "대학(원)생이라면", size=22, bold=True, color=OFF_BLACK)
multi(s, Inches(1.6), Inches(4.5), Inches(4.4), Inches(2.0), [
    "Google AI Pro 플랜 1년 무료",
    "학생 인증만 하면 끝",
    "Gemini Advanced + NotebookLM Plus",
    "",
    ("gemini.google/students", BLUE),
], size=14, color=MID_GRAY, spacing=5)

card(s, Inches(7.0), Inches(3.5), Inches(5.2), Inches(3.3), SUBTLE_BG)
deco_dot(s, Inches(7.4), Inches(3.9), BLUE_LIGHT)
txt(s, Inches(7.8), Inches(3.8), Inches(4), Inches(0.5),
    "누구나 가능", size=22, bold=True, color=OFF_BLACK)
multi(s, Inches(7.4), Inches(4.5), Inches(4.4), Inches(2.0), [
    "Coursera x Google AI 과정 수강 시",
    "Gemini 3개월 무료",
    "",
    ("7일 무료 체험 후 과금 주의", RED_SOFT),
], size=14, color=MID_GRAY, spacing=5)
page(s, 7, TOTAL)

# ═══════════════════════════════════════════
# 8 — 커스텀 인스트럭션
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_circle(s, W - Inches(4), Inches(-1.5), 5, RGBColor(0xF0, 0xF0, 0xF0))
deco_dot(s, ML, Inches(0.8), BLUE)

txt(s, ML, Inches(1.0), CW, Inches(1.5),
    "커스텀\n인스트럭션.", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(2.8), Inches(2))

txt(s, ML, Inches(3.5), Inches(10), Inches(1.0),
    '"나는 이런 사람이야,\n이렇게 답해줘"', size=28, color=DARK_GRAY, line_spacing=1.4)

txt(s, ML, Inches(4.8), Inches(10), Inches(0.5),
    "라고 AI에게 미리 설정해두는 것.", size=18, color=MID_GRAY)

y = Inches(5.6)
for text, is_accent in [
    ("매번 같은 말 반복할 필요 없음", False),
    ("내 스타일에 맞는 답변을 받을 수 있음", False),
    ("한 번 설정하면 모든 대화에 적용", True),
]:
    deco_dot(s, ML, y + Inches(0.07), BLUE if is_accent else RGBColor(0xDD, 0xDD, 0xDD))
    txt(s, Inches(1.7), y, Inches(8), Inches(0.4),
        text, size=17, bold=is_accent, color=BLUE if is_accent else MID_GRAY)
    y += Inches(0.5)
page(s, 8, TOTAL)

# ═══════════════════════════════════════════
# 9 — 환각 차단
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.8), BLUE)
txt(s, ML, Inches(1.0), CW, Inches(1.2),
    "같은 질문,\n완전히 다른 답변.", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(2.8), Inches(2))

txt(s, ML, Inches(3.3), Inches(10), Inches(0.5),
    '"2026년 한국 소상공인 종사자 수가 몇 명이야?"', size=16, color=MID_GRAY)

# Left — bad
card(s, ML, Inches(4.2), Inches(5.2), Inches(2.8), SUBTLE_BG)
txt(s, Inches(1.6), Inches(4.4), Inches(4), Inches(0.3),
    "인스트럭션 없음", size=12, bold=True, color=MID_GRAY)
line(s, Inches(1.6), Inches(4.85), Inches(1), RED_SOFT)
txt(s, Inches(1.6), Inches(5.1), Inches(4.4), Inches(1.0),
    "그럴듯한 숫자를\n자신 있게 말함", size=24, bold=True, color=RED_SOFT, line_spacing=1.3)

# Right — good
card(s, Inches(7.0), Inches(4.2), Inches(5.2), Inches(2.8), RGBColor(0xEA, 0xF0, 0xFF))
txt(s, Inches(7.4), Inches(4.4), Inches(4), Inches(0.3),
    "인스트럭션 있음", size=12, bold=True, color=BLUE)
line(s, Inches(7.4), Inches(4.85), Inches(1), BLUE)
txt(s, Inches(7.4), Inches(5.1), Inches(4.4), Inches(1.0),
    '"확실하지 않습니다,\n출처 확인 필요"', size=24, bold=True, color=BLUE_DEEP, line_spacing=1.3)
page(s, 9, TOTAL)

# ═══════════════════════════════════════════
# 10 — 3가지 타입
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.8), BLUE)
txt(s, ML, Inches(1.0), CW, Inches(0.8),
    "3가지 타입.", size=48, bold=True, color=OFF_BLACK)
line(s, ML, Inches(2.0), Inches(2))

cards_data = [
    ("01", "환각 차단형", "틀린 정보가\n싫은 사람", BLUE_DEEP),
    ("02", "글쓰기 특화형", "블로그, SNS\n뉴스레터 자주 쓰는 사람", BLUE),
    ("03", "업무 효율형", "보고서, 정리\n요약 자주 하는 사람", BLUE_LIGHT),
]
col_w = Inches(3.3)
gap = Inches(0.35)
for i, (num, title, desc, accent) in enumerate(cards_data):
    x = ML + (col_w + gap) * i
    card(s, x, Inches(2.7), col_w, Inches(3.8), SUBTLE_BG)
    txt(s, x + Inches(0.4), Inches(2.9), Inches(1.5), Inches(0.8),
        num, size=48, bold=True, color=GHOST_LIGHT)
    txt(s, x + Inches(0.4), Inches(3.8), col_w - Inches(0.8), Inches(0.5),
        title, size=22, bold=True, color=OFF_BLACK)
    line(s, x + Inches(0.4), Inches(4.45), Inches(1.2), accent)
    txt(s, x + Inches(0.4), Inches(4.7), col_w - Inches(0.8), Inches(1.2),
        desc, size=15, color=MID_GRAY, line_spacing=1.5)

txt(s, ML, Inches(6.8), Inches(10), Inches(0.4),
    "집에 가서 바로 복붙하세요 — 템플릿 배포합니다", size=16, bold=True, color=BLUE)
page(s, 10, TOTAL)

# ═══════════════════════════════════════════
# 11/12/13 — 인스트럭션 상세
# ═══════════════════════════════════════════
instr_slides = [
    ("환각 차단형", [
        '근거 없거나 불확실한 경우 "알 수 없습니다"라고 명시',
        '답변 전 단계별로 정보 검증, 모호한 부분은 "확실하지 않음" 표시',
        '확실한 정보만 사용하여 간결하게 답변',
        '추측이 불가피하면 "추측입니다"라고 밝히기',
        '모호한 질문이면 먼저 맥락을 요청',
        '출처/근거가 있으면 함께 제시',
    ]),
    ("글쓰기 특화형", [
        '불필요한 서론 없이 바로 본문으로 시작',
        '문어체가 아닌 자연스러운 구어체로 작성',
        '문장은 짧고 명확하게, 단락은 3줄 이내',
        '핵심 메시지는 첫 문장 또는 마지막 문장에 배치',
        '훅(hook) 문장을 첫 줄에 넣기',
        '이모지나 과도한 bullet point 사용 피하기',
    ]),
    ("업무 효율형", [
        '결론을 먼저 말하고, 근거는 그 다음에',
        '답변은 bullet point로 간결하게 정리',
        '인사말, 감탄사, 마무리 멘트 생략',
        '숫자나 데이터가 있으면 반드시 포함',
        '애매한 표현 대신 구체적이고 실행 가능한 표현',
        '모호하면 먼저 물어보기',
    ]),
]
for idx, (title, items) in enumerate(instr_slides):
    s = blank(); bg(s, SUBTLE_BG)
    deco_circle(s, W - Inches(3.5), Inches(-1), 4, RGBColor(0xE8, 0xE8, 0xE8))
    deco_dot(s, ML, Inches(0.6), BLUE)
    txt(s, Inches(1.65), Inches(0.5), Inches(5), Inches(0.3),
        "CUSTOM INSTRUCTION", size=11, bold=True, color=BLUE)
    txt(s, ML, Inches(1.0), CW, Inches(0.8),
        title, size=42, bold=True, color=OFF_BLACK)
    line(s, ML, Inches(2.0), Inches(2))

    y = Inches(2.7)
    for i, inst in enumerate(items):
        txt(s, ML, y, Inches(0.7), Inches(0.5),
            f"{i+1}", size=22, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        txt(s, Inches(2.3), y + Inches(0.03), Inches(9), Inches(0.5),
            inst, size=17, color=DARK_GRAY)
        if i < len(items) - 1:
            line(s, Inches(2.3), y + Inches(0.5), Inches(8), RGBColor(0xDD, 0xDD, 0xDD))
        y += Inches(0.7)
    page(s, 11 + idx, TOTAL)

# ═══════════════════════════════════════════
# 14 — Section ②
# ═══════════════════════════════════════════
s = blank(); bg(s, BG_BLACK)
ghost_num(s, "02", GHOST, right_align=True)
deco_circle(s, Inches(9), Inches(4), 5, RGBColor(0x12, 0x12, 0x12))
deco_dot(s, ML, Inches(1.2), BLUE)

txt(s, ML, Inches(1.4), Inches(2), Inches(0.4),
    "PART 02", size=13, bold=True, color=BLUE_LIGHT)
txt(s, ML, Inches(2.2), CW, Inches(2.5),
    "잘 된 거 뜯어서\n내 걸로 만들어라.", size=60, bold=True, color=WHITE, line_spacing=1.08)
line(s, ML, Inches(4.8), Inches(2.5))
txt(s, ML, Inches(5.2), Inches(10), Inches(0.8),
    "프롬프트 잘 짜려고 고민하지 마세요.\n구조는 훔치고 내용만 바꾸면 됩니다.",
    size=20, color=MID_GRAY, line_spacing=1.5)
page(s, 14, TOTAL)

# ═══════════════════════════════════════════
# 15 — 역설계 3단계
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.8), BLUE)
txt(s, ML, Inches(1.0), CW, Inches(1.2),
    "역설계\n3단계 프레임.", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(2.8), Inches(2))

steps = [
    ("01", "진단", "뭐가 잘 됐는지\n먼저 파악"),
    ("02", "역설계", "어떤 프롬프트였는지\n추출"),
    ("03", "검증 → 이식", "돌려보고\n내 변수로 교체"),
]
col_w = Inches(3.3)
gap = Inches(0.35)
for i, (num, name, desc) in enumerate(steps):
    x = ML + (col_w + gap) * i
    txt(s, x, Inches(3.5), col_w, Inches(1.0),
        num, size=72, bold=True, color=GHOST_LIGHT)
    txt(s, x, Inches(4.6), col_w, Inches(0.5),
        name, size=26, bold=True, color=OFF_BLACK)
    line(s, x, Inches(5.2), Inches(1.2), BLUE)
    txt(s, x, Inches(5.5), col_w, Inches(1.0),
        desc, size=15, color=MID_GRAY, line_spacing=1.5)

    # Arrow between columns
    if i < 2:
        txt(s, x + col_w + Inches(0.05), Inches(4.5), Inches(0.3), Inches(0.5),
            "→", size=24, color=BLUE, align=PP_ALIGN.CENTER)
page(s, 15, TOTAL)

# ═══════════════════════════════════════════
# 16 — 결과물 어디서?
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.8), BLUE)
txt(s, ML, Inches(1.0), CW, Inches(1.2),
    "어디서 잘 된\n결과물을 찾을까?", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(2.8), Inches(2))

# Two feature cards
for i, (name, url, desc, dot_c) in enumerate([
    ("지피테이블", "gptable.net", "국내 최대 프롬프트\n공유 커뮤니티", BLUE),
    ("PromptBase", "promptbase.com", "글로벌 최대 프롬프트 마켓\n결과물 예시 포함", BLUE_LIGHT),
]):
    x = ML + Inches(5.6) * i
    card(s, x, Inches(3.4), Inches(5.2), Inches(2.5), SUBTLE_BG)
    deco_dot(s, x + Inches(0.4), Inches(3.75), dot_c)
    txt(s, x + Inches(0.7), Inches(3.6), Inches(4), Inches(0.5),
        name, size=26, bold=True, color=OFF_BLACK)
    txt(s, x + Inches(0.4), Inches(4.2), Inches(4), Inches(0.3),
        url, size=13, color=dot_c)
    line(s, x + Inches(0.4), Inches(4.6), Inches(1), dot_c)
    txt(s, x + Inches(0.4), Inches(4.8), Inches(4.2), Inches(1.0),
        desc, size=14, color=MID_GRAY, line_spacing=1.5)

txt(s, ML, Inches(6.4), Inches(10), Inches(0.5),
    "블로그 글  ·  SNS 카드뉴스  ·  PPT  ·  보고서 — 어떤 결과물이든 역설계 가능",
    size=15, color=MID_GRAY)
page(s, 16, TOTAL)

# ═══════════════════════════════════════════
# 17 — 영어 프롬프트
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.8), BLUE)
txt(s, ML, Inches(1.0), CW, Inches(1.2),
    "영어 프롬프트가\n더 좋을까?", size=48, bold=True, color=OFF_BLACK, line_spacing=1.1)
line(s, ML, Inches(2.8), Inches(2))

table_clean(s, 4, 3, [
    ["", "영어 프롬프트", "한국어 프롬프트"],
    ["토큰 효율", "4.5배 높음", "상대적으로 비쌈"],
    ["추론 정확도", "복잡한 작업에서 약간 유리", "GPT-4o 이후 격차 줄어듦"],
    ["일반 업무", "차이 거의 없음", "차이 거의 없음"],
], ML, Inches(3.4), Inches(9))

txt(s, ML, Inches(6.0), Inches(10), Inches(0.8),
    "영어로 프롬프트 쓰고, 한국어로 답변 받기.\n성능도 챙기고 결과물도 한국어로.",
    size=22, bold=True, color=BLUE, line_spacing=1.4)
page(s, 17, TOTAL)

# ═══════════════════════════════════════════
# 18 — 시연 1
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.6), BLUE)
txt(s, Inches(1.65), Inches(0.5), Inches(5), Inches(0.3),
    "LIVE DEMO", size=11, bold=True, color=BLUE)
txt(s, ML, Inches(1.0), CW, Inches(0.8),
    "프롬프트 역설계", size=44, bold=True, color=OFF_BLACK)
line(s, ML, Inches(2.0), Inches(2))

flow = [
    ("01", "잘 된 결과물 하나 고르기"),
    ("02", "진단 — 왜 좋은지 분석"),
    ("03", "역설계 — 어떤 프롬프트인지 추출"),
    ("04", "검증 → 이식 — 내 변수로 교체"),
]
y = Inches(2.8)
for num, step in flow:
    deco_dot(s, ML, y + Inches(0.1), BLUE)
    txt(s, Inches(1.7), y, Inches(0.8), Inches(0.4),
        num, size=14, bold=True, color=BLUE_LIGHT)
    txt(s, Inches(2.5), y, Inches(8), Inches(0.4),
        step, size=22, color=DARK_GRAY)
    if num != "04":
        line(s, Inches(2.5), y + Inches(0.55), Inches(7), RGBColor(0xEE, 0xEE, 0xEE))
    y += Inches(0.85)

txt(s, ML, Inches(6.3), Inches(10), Inches(0.4),
    "이게 되면 어디든 똑같이 적용됩니다.", size=20, bold=True, color=BLUE)
page(s, 18, TOTAL)

# ═══════════════════════════════════════════
# 19 — Section ③
# ═══════════════════════════════════════════
s = blank(); bg(s, BG_BLACK)
ghost_num(s, "03", GHOST, right_align=True)
deco_circle(s, Inches(8), Inches(4), 5, RGBColor(0x12, 0x12, 0x12))
deco_dot(s, ML, Inches(1.2), BLUE)

txt(s, ML, Inches(1.4), Inches(2), Inches(0.4),
    "PART 03", size=13, bold=True, color=BLUE_LIGHT)
txt(s, ML, Inches(2.2), CW, Inches(2.5),
    "나만의 봇을\n만들어라.", size=60, bold=True, color=WHITE, line_spacing=1.08)
line(s, ML, Inches(4.8), Inches(2.5))
txt(s, ML, Inches(5.2), Inches(10), Inches(0.5),
    "매번 프롬프트 안 써도 됩니다", size=20, color=MID_GRAY)
page(s, 19, TOTAL)

# ═══════════════════════════════════════════
# 20 — 봇 만들기
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.6), BLUE)
txt(s, Inches(1.65), Inches(0.5), Inches(5), Inches(0.3),
    "LIVE DEMO", size=11, bold=True, color=BLUE)
txt(s, ML, Inches(1.0), CW, Inches(0.8),
    "나만의 봇 만들기", size=44, bold=True, color=OFF_BLACK)
line(s, ML, Inches(2.0), Inches(2))

y = Inches(2.8)
for text, is_accent in [
    ("GPTs (ChatGPT) 또는 Gems (Gemini) 라이브 세팅", False),
    ("역설계한 프롬프트를 봇에 그대로 심기", False),
    ("한 번 만들면 클릭 한 번으로 재사용", True),
]:
    deco_dot(s, ML, y + Inches(0.08), BLUE if is_accent else RGBColor(0xDD, 0xDD, 0xDD))
    txt(s, Inches(1.7), y, Inches(9), Inches(0.4),
        text, size=20, bold=is_accent, color=BLUE if is_accent else DARK_GRAY)
    y += Inches(0.6)

# comparison cards
card(s, ML, Inches(4.8), Inches(5.2), Inches(2.0), SUBTLE_BG)
txt(s, Inches(1.6), Inches(5.0), Inches(4), Inches(0.3),
    "매번 프롬프트 입력", size=13, bold=True, color=MID_GRAY)
line(s, Inches(1.6), Inches(5.4), Inches(1), LIGHT_GRAY)
txt(s, Inches(1.6), Inches(5.6), Inches(4), Inches(0.8),
    "복붙, 수정, 반복\n실수 가능", size=20, color=MID_GRAY, line_spacing=1.4)

card(s, Inches(7.0), Inches(4.8), Inches(5.2), Inches(2.0), BLUE_PALE)
txt(s, Inches(7.4), Inches(5.0), Inches(4), Inches(0.3),
    "봇으로 세팅", size=13, bold=True, color=BLUE)
line(s, Inches(7.4), Inches(5.4), Inches(1), BLUE)
txt(s, Inches(7.4), Inches(5.6), Inches(4), Inches(0.8),
    "클릭 한 번\n일관된 결과", size=20, bold=True, color=BLUE_DEEP, line_spacing=1.4)
page(s, 20, TOTAL)

# ═══════════════════════════════════════════
# 21 — 정리 + 숙제
# ═══════════════════════════════════════════
s = blank(); bg(s, WHITE)
deco_dot(s, ML, Inches(0.8), BLUE)
txt(s, ML, Inches(1.0), CW, Inches(0.8),
    "오늘 배운 것.", size=48, bold=True, color=OFF_BLACK)
line(s, ML, Inches(2.0), Inches(2))

items = [
    ("01", "AI한테 나를 기억시켜라", "커스텀 인스트럭션 설정"),
    ("02", "잘 된 거 뜯어서 내 걸로", "역설계 3단계 프레임"),
    ("03", "나만의 봇을 만들어라", "GPTs / Gems 활용"),
]
y = Inches(2.7)
for num, title, desc in items:
    deco_dot(s, ML, y + Inches(0.1), BLUE)
    txt(s, Inches(1.7), y, Inches(0.8), Inches(0.4),
        num, size=14, bold=True, color=BLUE_LIGHT)
    txt(s, Inches(2.5), y, Inches(4.5), Inches(0.4),
        title, size=20, bold=True, color=OFF_BLACK)
    txt(s, Inches(7.5), y, Inches(4.5), Inches(0.4),
        desc, size=16, color=MID_GRAY)
    if num != "03":
        line(s, Inches(2.5), y + Inches(0.55), Inches(8.5), RGBColor(0xEE, 0xEE, 0xEE))
    y += Inches(0.8)

# 숙제
card(s, ML, Inches(5.2), CW, Inches(1.8), SUBTLE_BG)
deco_dot(s, Inches(1.6), Inches(5.4), BLUE)
txt(s, Inches(2.0), Inches(5.35), Inches(3), Inches(0.3),
    "TODAY'S HOMEWORK", size=11, bold=True, color=BLUE)
txt(s, Inches(1.6), Inches(5.8), Inches(9.5), Inches(1.0),
    "집에 가서 내가 자주 쓰는 프롬프트 하나 골라서\n봇으로 만들어보세요.",
    size=22, bold=True, color=OFF_BLACK, line_spacing=1.4)
page(s, 21, TOTAL)

# ═══════════════════════════════════════════
# 22 — Q&A
# ═══════════════════════════════════════════
s = blank(); bg(s, BG_BLACK)
deco_circle(s, Inches(4), Inches(-2), 8, GHOST)
deco_circle(s, W - Inches(3), Inches(5), 4, RGBColor(0x14, 0x14, 0x14))
deco_dot(s, Inches(6.0), Inches(1.8), BLUE)
deco_dot(s, Inches(6.4), Inches(1.8), BLUE_LIGHT)
deco_dot(s, Inches(6.8), Inches(1.8), RGBColor(0x33, 0x33, 0x33))

txt(s, ML, Inches(2.2), CW, Inches(1.5),
    "Q&A", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
line(s, Inches(5.5), Inches(4.0), Inches(2.3))
txt(s, ML, Inches(4.4), CW, Inches(0.5),
    "심화 워크샵에 관심 있으시면 QR코드를 스캔해주세요",
    size=17, color=MID_GRAY, align=PP_ALIGN.CENTER)
txt(s, ML, Inches(5.8), CW, Inches(0.4),
    "감사합니다", size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)
page(s, 22, TOTAL)


# ── Save ──
prs.save("/Users/kim-yewon/Documents/airo-won/ai 강의/slides_bold.pptx")
print("Done! slides_bold.pptx created — 22 slides.")
